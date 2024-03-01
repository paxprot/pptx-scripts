import os
import sys
from pptx import Presentation

def process_text(input_text):
    translation_table = str.maketrans("—，。！？：；‘’“”", "-,.!?:;''\"\"")
    cleaned_text = input_text.translate(translation_table)

    # Remove leading "- " from each line
    cleaned_lines = [line.lstrip("- ").strip() for line in cleaned_text.splitlines()]
    non_empty_cleaned_lines = [line for line in cleaned_lines if line.strip()]
    return "\n".join(non_empty_cleaned_lines)

def extract_notes_from_pptx(pptx_file):
    # Load the PowerPoint presentation
    presentation = Presentation(pptx_file)

    # Create a directory for notes
    base_filename = os.path.splitext(os.path.basename(pptx_file))[0]
    notes_dir = os.path.join(os.path.dirname(pptx_file), f"{base_filename}_notes")
    os.makedirs(notes_dir, exist_ok=True)

    # Extract notes from each slide
    for slide_number, slide in enumerate(presentation.slides, start=1):
        notes = slide.notes_slide.notes_text_frame.text
        trimmed_notes = process_text(notes)

        # Save notes to a text file
        output_filename = os.path.join(notes_dir, f"Slide {slide_number}.txt")
        with open(output_filename, "w", encoding="utf-8") as f:
            f.write(trimmed_notes)

    print(f"Notes extracted and saved in directory: {notes_dir}")

if __name__ == "__main__":
    if len(sys.argv) != 2:
        print("Usage: python export_pptx_notes.py <input_presentation.pptx>")
        sys.exit(1)

    input_pptx = sys.argv[1]
    extract_notes_from_pptx(input_pptx)
