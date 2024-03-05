import os
import sys
from pptx import Presentation
from pptx.util import Inches
from lxml import etree

def autoplay_media(media):
    el_id = xpath(media.element, './/p:cNvPr')[0].attrib['id']
    el_cnt = xpath(
        media.element.getparent().getparent().getparent(),
        './/p:timing//p:video//p:spTgt[@spid="%s"]' % el_id,
    )[0]
    cond = xpath(el_cnt.getparent().getparent(), './/p:cond')[0]
    cond.set('delay', '0')

def xpath(el, query):
    nsmap = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}
    return etree.ElementBase.xpath(el, query, namespaces=nsmap)

def insert_audio_to_pptx(pptx_file):
    # Load the presentation
    presentation = Presentation(pptx_file)

    # Get the folder name for notes (e.g., "test_notes" for "test.pptx")
    audio_folder = os.path.join(os.path.dirname(pptx_file), "audio")

    # Iterate through slides
    for slide_number, slide in enumerate(presentation.slides, start=1):
        audio_file_path = os.path.join(audio_folder, f"Slide {slide_number}.mp3")

        # Check if audio file exists
        if os.path.exists(audio_file_path):
            audio_shape = slide.shapes.add_movie(
                audio_file_path,
                left=Inches(-1.5),
                top=0,
                width=Inches(1.5),
                height=Inches(1.5),
                poster_frame_image=None,
                mime_type='audio/mp3'
            )

            # Get this workaround to generate autoplay audio from:
            # https://github.com/scanny/python-pptx/issues/427#issuecomment-856724440
            autoplay_media(audio_shape)

    # Save the modified presentation
    output_pptx_file = os.path.splitext(pptx_file)[0] + "_with_audio.pptx"
    presentation.save(output_pptx_file)

    print(f"Audio files inserted into {output_pptx_file}")

if __name__ == "__main__":
    if len(sys.argv) != 2:
        print("Usage: python insert_audio.py <input_presentation.pptx>")
        sys.exit(1)

    input_pptx_file = sys.argv[1]
    insert_audio_to_pptx(input_pptx_file)
